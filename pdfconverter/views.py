import os
from io import BytesIO
from django.shortcuts import render
from django.conf import settings
from .forms import UploadForm
from docx import Document
from PIL import Image
from reportlab.pdfgen import canvas
from pptx import Presentation
import pandas as pd

def file_to_pdf(input_paths, output_path):

    pdf_canvas = canvas.Canvas(output_path)
    
    for input_path in input_paths:
        if input_path.endswith(".docx"):
            document = Document(input_path)

            for paragraph in document.paragraphs:
                pdf_canvas.drawString(100, 800, paragraph.text)
                pdf_canvas.showPage()

        elif input_path.endswith((".jpg", ".png")):
            pdf_canvas.drawImage(input_path, 0, 0)
            pdf_canvas.showPage()

        elif input_path.endswith(".pptx"):
            sunum = Presentation(input_path)

            for slide in sunum.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        pdf_canvas.drawString(100, 800, shape.text)
                        pdf_canvas.showPage()

        elif input_path.endswith((".xls", ".xlsx")):
            exc = pd.read_excel(input_path)
            pdf_canvas = canvas.Canvas(output_path)

            for row in exc.itertuples(index=False):
                pdf_canvas.drawString(50, 800, " | ".join(map(str, row)))
                pdf_canvas.showPage()

    pdf_canvas.save()


def convert_file(request):
    if request.method == "POST":
        form = UploadForm(request.POST, request.FILES)
        if form.is_valid():
            files = request.FILES.getlist('file')
            output_dir = os.path.join(settings.MEDIA_ROOT, "converted")
            os.makedirs(output_dir, exist_ok=True)

            input_paths = []
            for uploaded_file in files:
                input_path = os.path.join(settings.MEDIA_ROOT, uploaded_file.name)

                with open(input_path, 'wb') as f:
                    for chunk in uploaded_file.chunks():
                        f.write(chunk)

                input_paths.append(input_path)

            if len(input_paths) == 1:
                original_name = os.path.splitext(os.path.basename(input_paths[0]))[0]
                output_filename = f"{original_name}.pdf"
            else:
                output_filename = "birlesikdosya.pdf"

            output_path = os.path.join(output_dir, output_filename)

            try:
                file_to_pdf(input_paths, output_path)
                success_message = f"Dönüştürme Başarılı: {output_filename}"
            except Exception as e:
                success_message = f"Dönüştürme Hatalı: {str(e)}"

            return render(request, "pdfconverter/convert.html", {
                'form': form,
                'success': success_message
            })
    else:
        form = UploadForm()

    return render(request, "pdfconverter/convert.html", {"form": form})